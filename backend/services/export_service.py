from docx.shared import Pt
from io import BytesIO
from docx import Document
from docx.shared import Pt
from pptx import Presentation

from backend.models import Project, Section   # ✅ FIX
from backend.database import db              # optional but useful

def export_docx(project_id):
    project = Project.query.get(project_id)
    sections = Section.query.filter_by(project_id=project_id).all()

    doc = Document()

    # Title
    title = doc.add_heading(project.topic, 0)
    title.style.font.size = Pt(24)

    for sec in sections:
        # Section title
        h = doc.add_heading(sec.title, level=1)
        h.style.font.size = Pt(18)

        text = sec.refined_content or sec.content or ""

        # Split into paragraphs
        for line in text.split("\n"):
            line = line.strip()
            if not line:
                continue

            # Bullet point
            if line.startswith("- "):
                p = doc.add_paragraph(line[2:], style='List Bullet')
            else:
                p = doc.add_paragraph(line)

            p.style.font.size = Pt(12)

        doc.add_page_break()

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer
def export_pptx(project_id):
    project = Project.query.get(project_id)
    sections = Section.query.filter_by(project_id=project_id).all()

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = project.topic

    # Content slides
    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sec.title

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        text = sec.refined_content or sec.content or ""

        for line in text.split("\n"):
            line = line.strip()
            if not line:
                continue

            if line.startswith("- "):
                p = body.add_paragraph()
                p.text = line[2:]
                p.level = 1
            else:
                p = body.add_paragraph()
                p.text = line
                p.level = 0

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
